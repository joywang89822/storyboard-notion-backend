# -*- coding: utf-8 -*-
"""通用分鏡腳本 PPTX 產生器。

跟 build_storyboard_pptx.py（鐵支翻盤篇一次性版本）排版一致，但吃 JSON 輸入，
不寫死任何專案資料，讓 Notion 自動化後端可以重複呼叫。

JSON 結構見同資料夾 SCHEMA.md。

CLI 用法：
    python pptx_builder.py <script.json> <output.pptx> [--db-dir DIR] [--spec-img PATH]
"""
import sys
import os
import re
import json
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

FONT = "Microsoft JhengHei"
DARK = RGBColor(0x40, 0x40, 0x40)
LGREY = RGBColor(0xEC, 0xEC, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
RED = RGBColor(0xC0, 0x00, 0x00)
BORDER = RGBColor(0xB0, 0xB0, 0xB0)
GREY = RGBColor(0x90, 0x90, 0x90)
YELLOW = RGBColor(0xFF, 0xF2, 0x99)
CHECKBOX_CHARS = ("☑", "☐")
NO_STYLE_TABLE_GRID = "{5940675A-B579-460E-94D1-54222C63F5DA}"

# 固定欄位選項（跟 Notion「範本」頁面的 toggle 選項一致）。
# 用來把 meta 裡打勾的值，畫成「☑A　☐B　☑C」這種跟現在排版一樣的列。
FIELD_OPTIONS = {
    "遊戲": ["中文德撲", "印尼德撲", "印尼合併", "Wow slot"],
    "投放媒體": ["FB", "Google", "TikTok"],
    "單支秒數上限": ["6s", "15s", "30s", "45s"],
    "是否做長短秒": ["長秒", "短秒"],
    "素材類型": ["全新", "部分沿用舊素材", "舊素材微調", "relanguage", "resize"],
    "影像特殊需求": ["3D", "2D", "真人"],
    "音效特殊需求": ["人聲配音", "AI語音", "先配樂再做畫面", "需創作音樂"],
}

# 尺寸對照表（跟 Notion「素材尺寸需求」一致）：(類別, [(素材類型, 檔案格式, [尺寸...], 備註), ...])
SIZE_TAXONOMY = [
    ("適用廣告/社群素材", [
        ("Banner", "JPG", ["1920×1080", "1080×1920", "1080×1350", "1080×1080", "1200×628"], ""),
        ("Video", "mp4", ["1920×1080", "1080×1920", "1080×1350", "1080×1080"], ""),
        ("Gif", "JPG", ["1920×1080", "1080×1920", "1080×1350", "1080×1080", "1200×628"], ""),
        ("Gif", "mp4", ["1920×1080", "1080×1920", "1080×1350", "1080×1080"], ""),
    ]),
    ("適用商店頁素材", [
        ("Screenshot", "JPG",
         ["1920×1080（AN）", "2732×2048（iOS iPad）", "2688×1242", "2208×1242（iOS手機）",
          "1080×1920（AN）", "2048×2732（iOS iPad）", "1242×2688", "1242×2208"], ""),
        ("主題圖片", "JPG", ["1024×500"], ""),
        ("Liveops", "JPG", ["1080×1920", "1920×1080", "1080×1080"], ""),
        ("主題影片", "mp4", ["1920×1080", "1920×886"], ""),
        ("Icon", "JPG", ["512×512"], "<1MB"),
    ]),
]


def add_picture_fit(slide, path, left, top, max_w, max_h):
    """在 (max_w, max_h) 的框內完整顯示圖片，等比縮放、置中，絕不裁切。"""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    disp_w = int(iw * scale)
    disp_h = int(ih * scale)
    x = int(left + (max_w - disp_w) / 2)
    y = int(top + (max_h - disp_h) / 2)
    return slide.shapes.add_picture(path, x, y, width=disp_w, height=disp_h)


def txbox(slide, l, t, w, h, text, size=12, bold=False, color=BLACK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, fill=None, border=True,
          checkbox_red=False, red_terms=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.fill.solid() if fill else box.fill.background()
    if fill:
        box.fill.fore_color.rgb = fill
    if border:
        box.line.color.rgb = BORDER
        box.line.width = Pt(0.75)
    else:
        box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if checkbox_red and any(c in line for c in CHECKBOX_CHARS):
            for seg in re.split(r"([☑☐])", line):
                if not seg:
                    continue
                run = p.add_run()
                run.text = seg
                run.font.size = Pt(size)
                run.font.name = FONT
                if seg == "☑":
                    run.font.bold = True
                    run.font.color.rgb = RED
                else:
                    run.font.bold = bold
                    run.font.color.rgb = color
        elif red_terms:
            pattern = "(" + "|".join(re.escape(t) for t in red_terms) + ")"
            for seg in re.split(pattern, line):
                if not seg:
                    continue
                run = p.add_run()
                run.text = seg
                run.font.size = Pt(size)
                run.font.name = FONT
                if seg in red_terms:
                    run.font.bold = True
                    run.font.color.rgb = RED
                else:
                    run.font.bold = bold
                    run.font.color.rgb = color
        else:
            p.text = line
            for r in p.runs:
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.color.rgb = color
                r.font.name = FONT
    return box


def header_bar(slide, l, t, w, h, text, size=16):
    return txbox(slide, l, t, w, h, text, size=size, bold=True, color=WHITE,
                 fill=DARK, anchor=MSO_ANCHOR.MIDDLE)


def make_table(slide, l, t, w, h, n_rows, n_cols, col_widths):
    gtbl = slide.shapes.add_table(n_rows, n_cols, l, t, w, h).table
    gtbl.first_row = False
    gtbl.horz_banding = False
    style_id = gtbl._tbl.find(qn('a:tblPr')).find(qn('a:tableStyleId'))
    if style_id is not None:
        style_id.text = NO_STYLE_TABLE_GRID
    for i, cw in enumerate(col_widths):
        gtbl.columns[i].width = cw
    return gtbl


def set_cell(cell, text, size=11, bold=False, color=BLACK, fill=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    cell.vertical_anchor = anchor
    cell.margin_left = Pt(6); cell.margin_right = Pt(6)
    cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = FONT


def _checkbox_line(field_name, checked_values):
    options = FIELD_OPTIONS[field_name]
    checked = set(checked_values or [])
    parts = []
    for opt in options:
        mark = "☑" if opt in checked else "☐"
        parts.append(f"{mark}{opt}")
    extra = [v for v in checked if v not in options]
    if extra:
        parts.append("☑其他：" + "、".join(extra))
    else:
        parts.append("☐其他：")
    return "　".join(parts)


def _sum_seconds(shots):
    total = 0.0
    for sh in shots:
        m = re.match(r"([\d.]+)", str(sh.get("seconds", "0") or "0"))
        if m:
            total += float(m.group(1))
    if total == int(total):
        return f"{int(total)}s"
    return f"{total}s"


def _resolve_img(path, base_dir):
    if not path:
        return None
    p = path if os.path.isabs(path) else os.path.join(base_dir, path)
    return p if os.path.exists(p) else None


def build_pptx(data, out_path, base_dir=".", spec_img=None):
    """data: 見 SCHEMA.md 的 dict。out_path: 輸出 pptx 路徑。
    base_dir: shots[].image 相對路徑的基準資料夾。spec_img: 最後一頁素材規範的示意圖，可留空。"""
    meta = data.get("meta", {})
    scenes_in = data.get("scenes", [])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    # ============================================================
    # Slide 1: 基本資訊
    # ============================================================
    slide = prs.slides.add_slide(BLANK)
    L, W = Inches(0.4), Inches(12.5)
    LBL_W = Inches(2.0)
    row_h = Inches(0.48)
    y = Inches(0.4)

    asset_name = meta.get("素材名稱", "")
    total_count = meta.get("總支數", "")
    total_note = meta.get("總支數備註", "")
    total_text = str(total_count) if not total_note else f"{total_count}，備註：{total_note}"

    rows = [
        ("遊戲／素材名稱", asset_name),
        ("投放媒體", _checkbox_line("投放媒體", meta.get("投放媒體"))),
        ("總支數", total_text),
        ("單支秒數上限", _checkbox_line("單支秒數上限", meta.get("單支秒數上限"))),
        ("是否做長短秒", _checkbox_line("是否做長短秒", meta.get("是否做長短秒"))),
        ("素材類型", _checkbox_line("素材類型", meta.get("素材類型"))),
        ("影像特殊需求", _checkbox_line("影像特殊需求", meta.get("影像特殊需求"))),
        ("音效特殊需求", _checkbox_line("音效特殊需求", meta.get("音效特殊需求"))),
    ]
    for i, (label, val) in enumerate(rows):
        is_first = i == 0
        txbox(slide, L, y, LBL_W, row_h, label, size=16, bold=True, fill=LGREY)
        txbox(slide, L + LBL_W, y, W - LBL_W, row_h, val, size=12, bold=is_first, checkbox_red=True)
        y += row_h

    y += Inches(0.1)
    rows2 = [
        ("畫面風格", meta.get("畫面風格", "")),
        ("配樂風格", meta.get("配樂風格", "")),
        ("廣告TA", meta.get("廣告TA", "")),
        ("廣告目的", meta.get("廣告目的", "")),
        ("廣告鉤子", meta.get("廣告鉤子", "")),
    ]
    for label, val in rows2:
        txbox(slide, L, y, LBL_W, row_h, label, size=16, bold=True, fill=LGREY)
        txbox(slide, L + LBL_W, y, W - LBL_W, row_h, val, size=12)
        y += row_h

    # ============================================================
    # Slide 2: 素材尺寸
    # ============================================================
    slide = prs.slides.add_slide(BLANK)
    header_bar(slide, Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.6),
               "素材尺寸（需求標黃底，紅字優先對稿）", size=18)

    selected_sizes = meta.get("素材尺寸", {})  # {"Video": ["1920×1080", ...], ...}
    priority_sizes = set(meta.get("優先對稿尺寸") or [])

    col_x = [Inches(0.4), Inches(2.2), Inches(4.0), Inches(5.6), Inches(11.6)]
    col_w = [Inches(1.8), Inches(1.8), Inches(1.6), Inches(6.0), Inches(1.3)]
    hdr_y = Inches(1.0)
    hdrs = ["類別", "素材類型", "檔案格式", "尺寸", "備註"]
    for x, w, h in zip(col_x, col_w, hdrs):
        txbox(slide, x, hdr_y, w, Inches(0.45), h, size=12, bold=True, fill=LGREY)

    y = hdr_y + Inches(0.45)
    rh = Inches(0.47)
    for cat, items in SIZE_TAXONOMY:
        cat_h = rh * len(items)
        txbox(slide, col_x[0], y, col_w[0], cat_h, cat, size=10, bold=True)
        ty = y
        for typ, fmt, sizes, base_note in items:
            checked = [s for s in sizes if s in (selected_sizes.get(typ) or [])]
            is_required = bool(checked)
            fillc = YELLOW if is_required else WHITE
            note = "本次需求" if is_required and not base_note else base_note
            txbox(slide, col_x[1], ty, col_w[1], rh, typ, size=10)
            txbox(slide, col_x[2], ty, col_w[2], rh, fmt, size=10)
            txbox(slide, col_x[3], ty, col_w[3], rh, "、".join(sizes), size=10, color=BLACK,
                  fill=fillc, red_terms=[s for s in checked if s in priority_sizes] or None)
            txbox(slide, col_x[4], ty, col_w[4], rh, note, size=10)
            ty += rh
        y += cat_h

    # ============================================================
    # Normalize scenes
    # ============================================================
    scenes = []
    for s in scenes_in:
        shots = s.get("shots", [])
        scenes.append({
            "no": s.get("id"),
            "dur": s.get("seconds") or _sum_seconds(shots),
            "layoutObjects": s.get("layoutObjects", []),
            "music": s.get("music", ""),
            "sfx": s.get("sfx", ""),
            "note": s.get("note", ""),
            "shots": [
                {
                    "img": _resolve_img(sh.get("image"), base_dir),
                    "angle": sh.get("angle", ""),
                    "pos": sh.get("position", ""),
                    "action": sh.get("action", ""),
                    "caption": sh.get("line", ""),
                    "dur": sh.get("seconds", ""),
                    "note": sh.get("note", ""),
                }
                for sh in shots
            ],
        })

    # ============================================================
    # Slide 3: 分鏡總表
    # ============================================================
    if scenes:
        slide = prs.slides.add_slide(BLANK)
        header_bar(slide, Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.6), "分鏡總表", size=18)
        txbox(slide, Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.35),
              "僅動態素材（gif/video）才需填寫", size=11, bold=True, color=RED, fill=LGREY, border=False)

        n = len(scenes)
        first_col_w = Inches(1.1)
        col_w = Inches(11.4) // n
        sum_col_widths = [first_col_w] + [col_w] * n

        row_lbl_y = Inches(1.35)
        row_lbl_h = Inches(0.45)
        img_h = Inches(4.15)
        caption_h = Inches(0.6)
        sum_tbl_h = row_lbl_h + img_h + caption_h

        sum_tbl = make_table(slide, Inches(0.4), row_lbl_y, Inches(12.5), sum_tbl_h, 3, n + 1, sum_col_widths)
        sum_tbl.rows[0].height = row_lbl_h
        sum_tbl.rows[1].height = img_h
        sum_tbl.rows[2].height = caption_h

        set_cell(sum_tbl.cell(0, 0), "分鏡", size=12, bold=True, fill=LGREY)
        for i, s in enumerate(scenes):
            set_cell(sum_tbl.cell(0, i + 1), str(s["no"]), size=12, bold=True, fill=LGREY, align=PP_ALIGN.CENTER)

        set_cell(sum_tbl.cell(1, 0), "示意圖", size=12, bold=True, fill=LGREY)
        for i in range(n):
            set_cell(sum_tbl.cell(1, i + 1), "", size=10)

        set_cell(sum_tbl.cell(2, 0), "台詞／字卡", size=12, bold=True, fill=LGREY)
        for i, s in enumerate(scenes):
            caps = [sh["caption"] for sh in s["shots"] if sh["caption"]]
            set_cell(sum_tbl.cell(2, i + 1), "、".join(caps), size=11)

        img_y = row_lbl_y + row_lbl_h
        x = Inches(0.4) + first_col_w
        for s in scenes:
            thumb_h = img_h / max(len(s["shots"]), 1)
            ty = img_y
            for sh in s["shots"]:
                if sh["img"]:
                    add_picture_fit(slide, sh["img"], x + Inches(0.05), ty, col_w - Inches(0.1), thumb_h - Inches(0.05))
                else:
                    txbox(slide, x + Inches(0.05), ty, col_w - Inches(0.1), thumb_h - Inches(0.05),
                          "尚未提供\n參考素材", size=10, color=GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                ty += thumb_h
            x += col_w

    # ============================================================
    # Slides 4+: one per 分鏡
    # ============================================================
    for s in scenes:
        slide = prs.slides.add_slide(BLANK)

        bar_h = Inches(0.55)
        head_col_w = [Inches(1.0), Inches(0.8), Inches(1.0), Inches(1.0)]
        head_tbl_w = sum(head_col_w, Emu(0))
        head_tbl = make_table(slide, Inches(0.4), Inches(0.3), head_tbl_w, bar_h, 1, 4, head_col_w)
        set_cell(head_tbl.cell(0, 0), "分鏡", size=13, bold=True, color=WHITE, fill=DARK, align=PP_ALIGN.CENTER)
        set_cell(head_tbl.cell(0, 1), str(s["no"]), size=13, bold=True, color=WHITE, fill=RGBColor(0x70, 0x70, 0x70), align=PP_ALIGN.CENTER)
        set_cell(head_tbl.cell(0, 2), "秒數", size=13, bold=True, color=WHITE, fill=DARK, align=PP_ALIGN.CENTER)
        set_cell(head_tbl.cell(0, 3), s["dur"], size=13, bold=True, color=WHITE, fill=RGBColor(0x70, 0x70, 0x70), align=PP_ALIGN.CENTER)

        img_left = Inches(0.4)
        img_top = Inches(0.95)
        img_w = head_tbl_w
        img_area_h = Inches(6.15)
        n_shots = max(len(s["shots"]), 1)
        thumb_h = img_area_h / n_shots
        ty = img_top
        note_lines = [f"分鏡備註：⚠ {s['note']}"] if s.get("note") else []
        for i, sh in enumerate(s["shots"], start=1):
            if sh["img"]:
                add_picture_fit(slide, sh["img"], img_left, ty, img_w, thumb_h - Inches(0.1))
            else:
                txbox(slide, img_left, ty, img_w, thumb_h - Inches(0.1),
                      "尚未提供參考素材\n（資料庫沒有對應素材，或可在 Notion 這個鏡頭底下貼圖）",
                      size=10, color=GREY,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            txbox(slide, img_left, ty, Inches(0.55), Inches(0.32), f"圖{i}", size=11, bold=True,
                  color=WHITE, fill=DARK, border=False)
            if sh.get("note"):
                note_lines.append(f"圖{i}：⚠ {sh['note']}")
            ty += thumb_h

        info_l = Inches(4.35)
        rh2 = Inches(0.45)
        sub_h = Inches(0.35)
        row_h2 = Inches(0.65)
        col_widths = [Inches(2.2), Inches(4.0), Inches(1.5), Inches(0.85)]
        tbl_w = sum(col_widths, Emu(0))
        has_notes = bool(note_lines)
        note_row_h = Inches(0.28) * max(len(note_lines), 1)
        attr_rows = [(lo.get("label", ""), lo.get("desc", "")) for lo in s["layoutObjects"]]
        n_attr = len(attr_rows)
        n_rows = 1 + n_attr + 1 + n_shots + 2 + (1 if has_notes else 0)
        tbl_h = bar_h + rh2 * n_attr + sub_h + row_h2 * n_shots + rh2 * 2 + (note_row_h if has_notes else Inches(0))

        gtbl = make_table(slide, info_l, Inches(0.3), tbl_w, tbl_h, n_rows, 4, col_widths)

        r = 0
        gtbl.rows[r].height = bar_h
        set_cell(gtbl.cell(r, 0), "排版物件", size=12, bold=True, color=WHITE, fill=DARK, align=PP_ALIGN.CENTER)
        gtbl.cell(r, 1).merge(gtbl.cell(r, 3))
        set_cell(gtbl.cell(r, 1), "說明", size=12, bold=True, color=WHITE, fill=DARK, align=PP_ALIGN.CENTER)
        r += 1

        for label, val in attr_rows:
            gtbl.rows[r].height = rh2
            set_cell(gtbl.cell(r, 0), label, size=12, bold=True)
            gtbl.cell(r, 1).merge(gtbl.cell(r, 3))
            set_cell(gtbl.cell(r, 1), val, size=11)
            r += 1

        gtbl.rows[r].height = sub_h
        for c, h_ in enumerate(["視角/鏡位", "動態說明", "文案/對白", "秒數"]):
            set_cell(gtbl.cell(r, c), h_, size=12, bold=True, color=WHITE, fill=DARK, align=PP_ALIGN.CENTER)
        r += 1

        for sh in s["shots"]:
            gtbl.rows[r].height = row_h2
            angle_pos = f"{sh['angle']}／{sh['pos']}" if sh.get("pos") else sh["angle"]
            set_cell(gtbl.cell(r, 0), angle_pos, size=10)
            set_cell(gtbl.cell(r, 1), sh["action"], size=10)
            set_cell(gtbl.cell(r, 2), sh["caption"], size=10, bold=bool(sh["caption"]))
            set_cell(gtbl.cell(r, 3), sh["dur"], size=10, align=PP_ALIGN.CENTER)
            r += 1

        gtbl.rows[r].height = rh2
        set_cell(gtbl.cell(r, 0), "音樂", size=12, bold=True)
        gtbl.cell(r, 1).merge(gtbl.cell(r, 3))
        set_cell(gtbl.cell(r, 1), s["music"], size=11)
        r += 1
        gtbl.rows[r].height = rh2
        set_cell(gtbl.cell(r, 0), "音效", size=12, bold=True)
        gtbl.cell(r, 1).merge(gtbl.cell(r, 3))
        set_cell(gtbl.cell(r, 1), s["sfx"], size=11)
        r += 1

        if has_notes:
            gtbl.rows[r].height = note_row_h
            set_cell(gtbl.cell(r, 0), "備註", size=12, bold=True)
            gtbl.cell(r, 1).merge(gtbl.cell(r, 3))
            set_cell(gtbl.cell(r, 1), "\n".join(note_lines), size=10)

    # ============================================================
    # 最後一頁：素材規範（公司固定樣板，每份分鏡腳本都要附上，內容不隨專案異動）
    # ============================================================
    slide = prs.slides.add_slide(BLANK)
    header_bar(slide, Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.6), "素材規範", size=18)
    txbox(slide, Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.35),
          "統一適用後續每支廣告素材（商店頁不適用），並將不在腳本中特別註明",
          size=11, bold=True, color=RED, fill=LGREY, border=False)

    spec_col_x = [Inches(0.4), Inches(2.3), Inches(4.3), Inches(11.4)]
    spec_col_w = [Inches(1.9), Inches(2.0), Inches(7.1), Inches(1.5)]
    spec_hdr_y = Inches(1.35)
    spec_hdr_h = Inches(0.4)
    for x, w, h in zip(spec_col_x, spec_col_w, ["類別", "說明", "規格／示意圖", "備註"]):
        txbox(slide, x, spec_hdr_y, w, spec_hdr_h, h, size=12, bold=True, fill=LGREY)

    spec_rows = [
        ("尺寸安全框", "重點訊息請優先\n放置於安全區內", None, "", Inches(2.1)),
        ("logo", "畫面左上置頂",
         "依素材內容套用，若背景太花logo不明顯，加背景相近色的色塊底", "", Inches(0.7)),
        ("分級標誌", "右下置底", "", "適用台灣、印尼遊戲", Inches(0.55)),
        ("警語", "下方置底\n最小化字級，分行最多不超過兩行\n依素材內容套用 (1)文字與背景相近色、(2)文字與文案相近字型、(3)可不加色塊底/加背景相近色的色塊底",
         "台灣：這不是一款賭博遊戲，且無任何方式可以贏得金錢獎勵或具有金錢價值的獎品。嚴禁賭博。\n"
         "印尼：Game ini tidak menyediakan aktivitas perjudian apa pun, chips dalam game tidak dapat ditukar dengan uang tunai atau barang nyata lainnya.",
         "適用印尼及中文德撲", Inches(1.7)),
    ]
    y = spec_hdr_y + spec_hdr_h
    for cat, desc, spec, note, h in spec_rows:
        txbox(slide, spec_col_x[0], y, spec_col_w[0], h, cat, size=11, bold=True)
        txbox(slide, spec_col_x[1], y, spec_col_w[1], h, desc, size=10)
        if spec is None and spec_img and os.path.exists(spec_img):
            add_picture_fit(slide, spec_img, spec_col_x[2], y, spec_col_w[2], h)
        else:
            txbox(slide, spec_col_x[2], y, spec_col_w[2], h, spec or "", size=10, color=BLACK)
        txbox(slide, spec_col_x[3], y, spec_col_w[3], h, note, size=10)
        y += h

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    return out_path


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("json_path")
    ap.add_argument("output_path")
    ap.add_argument("--db-dir", default=None, help="shots[].image 相對路徑的基準資料夾，預設為 json 檔所在資料夾")
    ap.add_argument("--spec-img", default=None, help="最後一頁素材規範示意圖路徑")
    args = ap.parse_args()

    with open(args.json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    base_dir = args.db_dir or os.path.dirname(os.path.abspath(args.json_path))
    out = build_pptx(data, args.output_path, base_dir=base_dir, spec_img=args.spec_img)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
